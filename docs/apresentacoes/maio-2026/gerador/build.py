# -*- coding: utf-8 -*-
import copy
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn

import theme as T
import data as D
from theme import (C, txt, rect, rounded, solid, set_border, no_line, no_shadow,
                   header, footer, canvas_bg, soft_shadow, shades)

prs = Presentation()
prs.slide_width = T.SLIDE_W
prs.slide_height = T.SLIDE_H
BLANK = prs.slide_layouts[6]
LM = Inches(0.62)
CW = Inches(12.09)  # largura útil


def add():
    return prs.slides.add_slide(BLANK)


# ---------- helpers de gráfico ----------
def _chart_no_border(chart):
    cs = chart._chartSpace
    spPr = cs.find(qn('c:spPr'))
    if spPr is None:
        spPr = cs.makeelement(qn('c:spPr'), {})
        cs.append(spPr)
    # noFill + no line
    for tag in ('a:noFill', 'a:solidFill', 'a:ln'):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    spPr.append(spPr.makeelement(qn('a:noFill'), {}))
    ln = spPr.makeelement(qn('a:ln'), {})
    ln.append(ln.makeelement(qn('a:noFill'), {}))
    spPr.append(ln)
    # fonte padrão do gráfico
    txpr = cs.find(qn('c:txPr'))
    if txpr is None:
        txpr = cs.makeelement(qn('c:txPr'), {})
        bodyPr = txpr.makeelement(qn('a:bodyPr'), {})
        lst = txpr.makeelement(qn('a:lstStyle'), {})
        p = txpr.makeelement(qn('a:p'), {})
        pPr = p.makeelement(qn('a:pPr'), {})
        defR = pPr.makeelement(qn('a:defRPr'), {'sz': '1100'})
        latin = defR.makeelement(qn('a:latin'), {'typeface': T.F_BODY})
        defR.append(latin)
        pPr.append(defR)
        p.append(pPr)
        txpr.append(bodyPr); txpr.append(lst); txpr.append(p)
        cs.append(txpr)


def _set_donut_hole(chart, pct):
    de = chart._chartSpace.find(qn('c:chart')).find(qn('c:plotArea')).find(qn('c:doughnutChart'))
    if de is not None:
        hs = de.find(qn('c:holeSize'))
        if hs is None:
            hs = de.makeelement(qn('c:holeSize'), {})
            de.append(hs)
        hs.set('val', str(pct))


def _point_colors(series, colors):
    for i, pt in enumerate(series.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        pt.format.line.color.rgb = T.WHITE
        pt.format.line.width = Pt(2.0)


def donut(slide, l, t, w, h, names, values, colors, hole=62, label_pct=True,
          label_color=T.WHITE, label_size=12):
    cd = CategoryChartData()
    cd.categories = names
    cd.add_series('S', values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, l, t, w, h, cd)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = False
    _chart_no_border(ch)
    _set_donut_hole(ch, hole)
    plot = ch.plots[0]
    ser = plot.series[0]
    _point_colors(ser, colors)
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_percentage = label_pct
    dl.show_value = not label_pct
    dl.number_format = '0%'
    dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.CENTER
    dl.font.size = Pt(label_size)
    dl.font.bold = True
    dl.font.name = T.F_MONO
    dl.font.color.rgb = label_color
    return ch


def hbar(slide, l, t, w, h, names, values, color, label_color=T.INK,
         max_scale=None):
    cd = CategoryChartData()
    cd.categories = names
    cd.add_series('S', values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, l, t, w, h, cd)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = False
    _chart_no_border(ch)
    plot = ch.plots[0]
    plot.gap_width = 55
    ser = plot.series[0]
    # degradê por tom: maior valor = cor mais forte (categorias vêm em ordem
    # crescente, então o tom mais cheio fica no último ponto / topo do gráfico)
    cols = list(reversed(shades('%02X%02X%02X' % (color[0], color[1], color[2]),
                                len(values))))
    _point_colors_bar(ser, cols)
    # rótulos de dados
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    dl.number_format = '#,##0'
    dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.size = Pt(11)
    dl.font.bold = True
    dl.font.name = T.F_MONO
    dl.font.color.rgb = label_color
    # eixos
    cat = ch.category_axis
    cat.has_major_gridlines = False
    cat.has_minor_gridlines = False
    cat.tick_labels.font.size = Pt(11)
    cat.tick_labels.font.name = T.F_BODY
    cat.tick_labels.font.color.rgb = T.BODY
    cat.format.line.color.rgb = T.HAIR2
    cat.format.line.width = Pt(0.75)
    val = ch.value_axis
    val.has_major_gridlines = False
    val.has_minor_gridlines = False
    val.visible = False
    if max_scale:
        val.maximum_scale = max_scale
        val.minimum_scale = 0
    # esconder linha do eixo de valor
    _axis_no_line(val)
    _axis_delete_labels(val)
    return ch


def _point_colors_bar(series, colors):
    for i, pt in enumerate(series.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        pt.format.line.fill.background()


def _axis_no_line(axis):
    spPr = axis._element.find(qn('c:spPr'))
    if spPr is None:
        spPr = axis._element.makeelement(qn('c:spPr'), {})
        axis._element.append(spPr)
    ln = spPr.makeelement(qn('a:ln'), {})
    ln.append(ln.makeelement(qn('a:noFill'), {}))
    spPr.append(ln)


def _axis_delete_labels(axis):
    el = axis._element
    tl = el.find(qn('c:majorTickMark'))
    # set tick + labels none
    for tag, val in (('c:majorTickMark', 'none'), ('c:minorTickMark', 'none'),
                     ('c:tickLblPos', 'none')):
        e = el.find(qn(tag))
        if e is None:
            e = el.makeelement(qn(tag), {})
            el.append(e)
        e.set('val', val)


# ---------- legenda / KPI custom ----------
def kpi_card(slide, l, t, w, h, label, value, sub, accent, value_color=None):
    card = rounded(slide, l, t, w, h, T.PANEL, radius=0.08)
    set_border(card, T.HAIR, 1.0)
    soft_shadow(card)
    rect(slide, l, t + Inches(0.16), Inches(0.07), h - Inches(0.32), accent)
    pad = l + Inches(0.28)
    txt(slide, pad, t + Inches(0.18), w - Inches(0.4), Inches(0.3),
        [{"text": label.upper(), "size": 10.5, "bold": True, "color": T.MUTED,
          "font": T.F_HEAD, "spacing": 1.4}])
    txt(slide, pad, t + Inches(0.46), w - Inches(0.4), Inches(0.6),
        [{"text": value, "size": 33, "bold": True,
          "color": value_color or T.INK, "font": T.F_MONO}])
    if sub:
        txt(slide, pad, t + h - Inches(0.42), w - Inches(0.4), Inches(0.3),
            [{"text": sub, "size": 10.5, "color": T.BODY, "font": T.F_BODY}])
    return card


def legend_row(slide, l, t, w, color, name, value, pct=None):
    rounded(slide, l, t + Inches(0.04), Inches(0.16), Inches(0.16), color, radius=0.3)
    runs = [{"text": name, "size": 12.5, "color": T.INK, "font": T.F_BODY, "bold": True}]
    txt(slide, l + Inches(0.30), t, Inches(3.4), Inches(0.3), [{"runs": runs}])
    vtxt = [{"text": value, "size": 12.5, "bold": True, "color": T.INK, "font": T.F_MONO,
             "align": PP_ALIGN.RIGHT}]
    if pct:
        vtxt.append({"text": "  " + pct, "size": 11, "color": T.MUTED, "font": T.F_MONO})
    txt(slide, l + Inches(3.0), t, w - Inches(3.0), Inches(0.3),
        [{"runs": vtxt, "align": PP_ALIGN.RIGHT}], align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1 — CAPA
# ============================================================
def slide_capa():
    s = add()
    canvas_bg(s, T.INK)
    # painel de acento à direita
    rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), T.INK)
    # bloco vertical teal à esquerda
    rect(s, Inches(0), Inches(0), Inches(0.42), Inches(7.5), T.BRAND)
    # faixa decorativa sutil
    band = rect(s, Inches(9.7), Inches(0), Inches(3.633), Inches(7.5), T.BRAND_D)
    # diagonal accent
    tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(11.0), Inches(4.9),
                             Inches(2.33), Inches(2.6))
    solid(tri, T.BRAND)
    tri.rotation = 0
    # marca
    txt(s, Inches(1.0), Inches(1.15), Inches(8), Inches(0.5),
        [{"runs": [
            {"text": "CLEAN", "size": 16, "bold": True, "color": T.WHITE,
             "font": T.F_HEAD, "spacing": 4.0},
            {"text": "TABACO", "size": 16, "bold": True, "color": C("9FC9CD"),
             "font": T.F_HEAD, "spacing": 4.0},
        ]}])
    rect(s, Inches(1.0), Inches(1.7), Inches(0.85), Pt(2.4), T.BRAND)
    # título principal
    txt(s, Inches(0.96), Inches(2.55), Inches(8.6), Inches(2.2),
        [
            {"text": "Relatório de", "size": 30, "color": C("BFD6D9"),
             "font": T.F_LIGHT, "line_spacing": 1.0, "space_after": 2},
            {"text": "Vendas & Produção", "size": 49, "bold": True, "color": T.WHITE,
             "font": T.F_HEAD, "line_spacing": 1.0},
        ])
    # data destaque
    txt(s, Inches(1.0), Inches(4.75), Inches(7), Inches(0.7),
        [{"runs": [
            {"text": "MAIO ", "size": 30, "bold": True, "color": T.BRAND,
             "font": T.F_HEAD, "spacing": 2.0},
            {"text": "2026", "size": 30, "color": C("9FC9CD"), "font": T.F_LIGHT,
             "spacing": 2.0},
        ]}])
    # mini KPIs no rodapé da capa
    kpis = [("VENDA TOTAL", D.br_int(D.VENDA_TOTAL)),
            ("PRODUÇÃO", D.br_int(D.PRODUCAO)),
            ("ESTOQUE FINAL", D.br_int(D.ESTOQUE_TOTAL[2]))]
    x = Inches(1.0)
    for lab, val in kpis:
        txt(s, x, Inches(5.95), Inches(2.3), Inches(0.9),
            [
                {"text": val, "size": 28, "bold": True, "color": T.WHITE,
                 "font": T.F_MONO, "space_after": 0},
                {"text": lab, "size": 10, "color": C("9FC9CD"), "font": T.F_BODY,
                 "spacing": 1.5},
            ])
        x += Inches(2.55)


# ============================================================
# SLIDE 2 — VENDA IN x VENDA OUT
# ============================================================
def slide_in_out():
    s = add()
    canvas_bg(s, T.WHITE)
    header(s, "Comercial · Maio 2026", "Venda IN  ×  Venda OUT", T.GREEN)
    total = D.VENDA_IN + D.VENDA_OUT
    donut(s, Inches(0.7), Inches(1.85), Inches(5.0), Inches(4.7),
          ["Venda IN", "Venda OUT"], [D.VENDA_IN, D.VENDA_OUT],
          [T.GREEN, T.AMBER], hole=60)
    # centro do donut: total
    txt(s, Inches(2.15), Inches(3.75), Inches(2.1), Inches(1.0),
        [
            {"text": D.br_int(total), "size": 30, "bold": True, "color": T.INK,
             "font": T.F_MONO, "align": PP_ALIGN.CENTER, "space_after": 0},
            {"text": "CAIXAS", "size": 10, "color": T.MUTED, "font": T.F_BODY,
             "align": PP_ALIGN.CENTER, "spacing": 1.5},
        ], align=PP_ALIGN.CENTER)
    # cartões KPI à direita
    kpi_card(s, Inches(6.4), Inches(2.0), Inches(6.3), Inches(1.65),
             "Venda IN  ·  78%", D.br_int(D.VENDA_IN),
             "Volume já faturado no mês", T.GREEN, T.GREEN)
    kpi_card(s, Inches(6.4), Inches(3.85), Inches(6.3), Inches(1.65),
             "Venda OUT  ·  22%", D.br_int(D.VENDA_OUT),
             "Saída de produtos no período", T.AMBER, T.AMBER)
    txt(s, Inches(6.4), Inches(5.75), Inches(6.3), Inches(0.5),
        [{"runs": [
            {"text": "Relação IN/OUT  ", "size": 11, "color": T.MUTED, "font": T.F_BODY},
            {"text": "3,5×", "size": 13, "bold": True, "color": T.INK, "font": T.F_MONO},
            {"text": "   —   a venda IN supera a OUT em 562 caixas", "size": 11,
             "color": T.MUTED, "font": T.F_BODY},
        ]}])
    footer(s, 2)


# ============================================================
# SLIDE 3 — IN / OUT / A FATURAR
# ============================================================
def slide_in_out_faturar():
    s = add()
    canvas_bg(s, T.WHITE)
    header(s, "Comercial · Maio 2026", "Venda IN · Venda OUT · A Faturar", T.BRAND)
    total = D.VENDA_TOTAL
    donut(s, Inches(0.7), Inches(1.85), Inches(5.0), Inches(4.7),
          ["Venda IN", "Venda OUT", "A faturar"],
          [D.VENDA_IN, D.VENDA_OUT, D.A_FATURAR],
          [T.GREEN, T.AMBER, T.GOLD], hole=60)
    txt(s, Inches(2.15), Inches(3.75), Inches(2.1), Inches(1.0),
        [
            {"text": D.br_int(total), "size": 28, "bold": True, "color": T.INK,
             "font": T.F_MONO, "align": PP_ALIGN.CENTER, "space_after": 0},
            {"text": "TOTAL", "size": 10, "color": T.MUTED, "font": T.F_BODY,
             "align": PP_ALIGN.CENTER, "spacing": 1.5},
        ], align=PP_ALIGN.CENTER)
    # legenda detalhada à direita
    rows = [
        (T.GREEN, "Venda IN", D.VENDA_IN, "68%"),
        (T.AMBER, "Venda OUT", D.VENDA_OUT, "19%"),
        (T.GOLD, "Venda IN a faturar (próx. mês)", D.A_FATURAR, "13%"),
    ]
    card = rounded(s, Inches(6.4), Inches(2.0), Inches(6.3), Inches(3.5), T.PANEL, 0.05)
    set_border(card, T.HAIR, 1.0)
    soft_shadow(card)
    y = Inches(2.45)
    for col, name, val, pct in rows:
        legend_row(s, Inches(6.75), y, Inches(5.6), col, name, D.br_int(val), pct)
        rect(s, Inches(6.75), y + Inches(0.42), Inches(5.6), Pt(0.8), T.HAIR)
        y += Inches(0.72)
    # total
    legend_row(s, Inches(6.75), y + Inches(0.05), Inches(5.6), T.INK, "Total geral",
               D.br_int(total), "100%")
    txt(s, Inches(6.4), Inches(5.75), Inches(6.3), Inches(0.8),
        [{"text": "Inclui 152 caixas de venda IN previstas para faturamento "
                  "no próximo mês.", "size": 11, "color": T.MUTED,
          "font": T.F_BODY, "line_spacing": 1.1}])
    footer(s, 3)


# ============================================================
# SLIDE 4 — VENDA TOTAL x PRODUÇÃO
# ============================================================
def slide_venda_producao():
    s = add()
    canvas_bg(s, T.WHITE)
    header(s, "Desempenho · Maio 2026", "Venda Total  ×  Produção", T.BLUE)
    donut(s, Inches(0.7), Inches(1.85), Inches(5.0), Inches(4.7),
          ["Produção", "Venda"], [D.PRODUCAO, D.VENDA_TOTAL],
          [T.BLUE, T.GREEN], hole=60)
    txt(s, Inches(2.15), Inches(3.78), Inches(2.1), Inches(1.0),
        [
            {"text": D.br_int(D.PRODUCAO + D.VENDA_TOTAL), "size": 24, "bold": True,
             "color": T.INK, "font": T.F_MONO, "align": PP_ALIGN.CENTER, "space_after": 0},
            {"text": "SOMA", "size": 9, "color": T.MUTED, "font": T.F_BODY,
             "align": PP_ALIGN.CENTER, "spacing": 1.5},
        ], align=PP_ALIGN.CENTER)
    kpi_card(s, Inches(6.4), Inches(2.0), Inches(6.3), Inches(1.65),
             "Produção  ·  55%", D.br_int(D.PRODUCAO),
             "Total produzido no mês", T.BLUE, T.BLUE)
    kpi_card(s, Inches(6.4), Inches(3.85), Inches(6.3), Inches(1.65),
             "Venda total  ·  45%", D.br_int(D.VENDA_TOTAL),
             "IN + OUT + a faturar (786 + 224 + 152)", T.GREEN, T.GREEN)
    txt(s, Inches(6.4), Inches(5.75), Inches(6.3), Inches(0.6),
        [{"runs": [
            {"text": "Saldo produção − venda  ", "size": 11, "color": T.MUTED,
             "font": T.F_BODY},
            {"text": "+" + D.br_int(D.PRODUCAO - D.VENDA_TOTAL), "size": 13,
             "bold": True, "color": T.BLUE, "font": T.F_MONO},
            {"text": "  caixas para estoque", "size": 11, "color": T.MUTED,
             "font": T.F_BODY},
        ]}])
    footer(s, 4)


# ============================================================
# SLIDES 5/6/7 — POR PRODUTO (barras)
# ============================================================
def slide_por_produto(idx_val, titulo, kicker, color, page):
    s = add()
    canvas_bg(s, T.WHITE)
    header(s, kicker, titulo, color)
    # ordena crescente para que o líder fique no TOPO do gráfico de barras
    rows = sorted(D.PRODUTOS_MAIO, key=lambda r: r[idx_val])
    names = [r[0] for r in rows]
    vals = [r[idx_val] for r in rows]
    total = sum(vals)
    ativos = sum(1 for v in vals if v > 0)
    lider = max(D.PRODUTOS_MAIO, key=lambda r: r[idx_val])
    # cartão do gráfico
    card = rounded(s, LM, Inches(1.78), Inches(8.55), Inches(5.05), T.PANEL, 0.03)
    set_border(card, T.HAIR, 1.0)
    soft_shadow(card)
    mx = max(vals) * 1.18
    hbar(s, Inches(0.78), Inches(1.95), Inches(8.25), Inches(4.7),
         names, vals, color, max_scale=mx)
    # painel lateral de destaques
    px = Inches(9.45)
    pw = Inches(3.26)
    kpi_card(s, px, Inches(1.78), pw, Inches(1.5), "Total do mês",
             D.br_int(total), "caixas", color, color)
    kpi_card(s, px, Inches(3.4), pw, Inches(1.5), "Líder", D.br_int(lider[idx_val]),
             lider[0], color, T.INK)
    kpi_card(s, px, Inches(5.02), pw, Inches(1.5), "SKUs com venda" if idx_val != 3
             else "SKUs produzidos", f"{ativos}/9", "produtos ativos", color, T.INK)
    footer(s, page)


# ============================================================
# TABELAS
# ============================================================
def _cell(cell, text, *, size=11, bold=False, color=T.BODY, font=T.F_BODY,
          align=PP_ALIGN.LEFT, fill=None, anchor=MSO_ANCHOR.MIDDLE):
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = anchor
    if fill is None:
        cell.fill.background()
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color


def _strip_table_style(tbl_graphic_frame):
    """Remove o estilo de tabela padrão (banding azulado)."""
    tbl = tbl_graphic_frame._element.graphic.graphicData.tbl
    tblPr = tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')
        for st in tblPr.findall(qn('a:tableStyleId')):
            tblPr.remove(st)


def _no_borders_then_custom(table):
    pass


def slide_estoque():
    s = add()
    canvas_bg(s, T.WHITE)
    header(s, "Estoque · Maio 2026", "Estoque Final — Matriz e Filial", T.SLATE)
    rows = D.ESTOQUE
    nrows = len(rows) + 2  # header + dados + total
    ncols = 5  # marca, matriz, filial, total, barra(participação)
    top = Inches(1.85)
    left = LM
    width = CW
    height = Inches(4.9)
    gf = s.shapes.add_table(nrows, ncols, left, top, width, height)
    _strip_table_style(gf)
    table = gf.table
    table.columns[0].width = Inches(5.0)
    table.columns[1].width = Inches(1.9)
    table.columns[2].width = Inches(1.9)
    table.columns[3].width = Inches(1.9)
    table.columns[4].width = Inches(1.39)
    # header
    hdrs = ["MARCA", "MATRIZ", "FILIAL", "TOTAL", "PART. %"]
    aligns = [PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT]
    for c in range(ncols):
        _cell(table.cell(0, c), hdrs[c], size=11, bold=True, color=T.WHITE,
              font=T.F_HEAD, align=aligns[c], fill=T.SLATE)
    table.rows[0].height = Inches(0.42)
    tot_total = D.ESTOQUE_TOTAL[2]
    for i, (marca, mz, fl, tt) in enumerate(rows):
        r = i + 1
        bg = T.WHITE if i % 2 == 0 else T.TINT
        _cell(table.cell(r, 0), marca, size=11.5, bold=True, color=T.INK,
              align=PP_ALIGN.LEFT, fill=bg)
        _cell(table.cell(r, 1), D.br_int(mz), size=11.5, color=T.BODY,
              font=T.F_MONO, align=PP_ALIGN.RIGHT, fill=bg)
        _cell(table.cell(r, 2), D.br_int(fl), size=11.5, color=T.BODY,
              font=T.F_MONO, align=PP_ALIGN.RIGHT, fill=bg)
        _cell(table.cell(r, 3), D.br_int(tt), size=11.5, bold=True, color=T.INK,
              font=T.F_MONO, align=PP_ALIGN.RIGHT, fill=bg)
        pct = (tt / tot_total * 100) if tot_total else 0
        _cell(table.cell(r, 4), f"{pct:.1f}%".replace('.', ','), size=10.5,
              color=T.MUTED, font=T.F_MONO, align=PP_ALIGN.RIGHT, fill=bg)
        table.rows[r].height = Inches(0.40)
    # total
    rr = nrows - 1
    _cell(table.cell(rr, 0), "TOTAL GERAL", size=11.5, bold=True, color=T.WHITE,
          font=T.F_HEAD, fill=T.INK)
    for c, v in enumerate(D.ESTOQUE_TOTAL, start=1):
        _cell(table.cell(rr, c), D.br_int(v), size=12, bold=True, color=T.WHITE,
              font=T.F_MONO, align=PP_ALIGN.RIGHT, fill=T.INK)
    _cell(table.cell(rr, 4), "100%", size=11, bold=True, color=T.WHITE,
          font=T.F_MONO, align=PP_ALIGN.RIGHT, fill=T.INK)
    table.rows[rr].height = Inches(0.46)
    _table_borders(table, ncols, nrows)
    footer(s, 8)


def slide_mensal(titulo, kicker, dados, total_row, color, page, total_label):
    s = add()
    canvas_bg(s, T.WHITE)
    header(s, kicker, titulo, color)
    meses = D.MESES
    headers = ["CÓD.", "PRODUTO"] + meses + ["CAIXAS", "CARTEIRA (R$)"]
    ncols = len(headers)  # 2 + 5 + 2 = 9
    nrows = len(dados) + 2
    top = Inches(1.78)
    gf = s.shapes.add_table(nrows, ncols, LM, top, CW, Inches(4.95))
    _strip_table_style(gf)
    table = gf.table
    widths = [1.05, 3.55, 0.78, 0.78, 0.78, 0.78, 0.78, 1.25, 1.56]
    for i, wv in enumerate(widths):
        table.columns[i].width = Inches(wv)
    aligns = [PP_ALIGN.CENTER, PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 5 + \
             [PP_ALIGN.RIGHT, PP_ALIGN.RIGHT]
    for c in range(ncols):
        _cell(table.cell(0, c), headers[c], size=10, bold=True, color=T.WHITE,
              font=T.F_HEAD, align=aligns[c], fill=color)
    table.rows[0].height = Inches(0.40)
    for i, (cod, prod, mvals, caixa, carteira) in enumerate(dados):
        r = i + 1
        bg = T.WHITE if i % 2 == 0 else T.TINT
        _cell(table.cell(r, 0), cod, size=8.5, color=T.MUTED, font=T.F_MONO,
              align=PP_ALIGN.CENTER, fill=bg)
        _cell(table.cell(r, 1), prod, size=10.5, bold=True, color=T.INK,
              align=PP_ALIGN.LEFT, fill=bg)
        for j, mv in enumerate(mvals):
            is_mai = (j == len(mvals) - 1)
            _cell(table.cell(r, 2 + j), D.br_cell(mv), size=10.5,
                  bold=is_mai, color=T.INK if is_mai else T.BODY, font=T.F_MONO,
                  align=PP_ALIGN.RIGHT, fill=(T.BRAND_L if is_mai else bg))
        _cell(table.cell(r, 7), D.br_dec(caixa), size=10.5, bold=True, color=T.INK,
              font=T.F_MONO, align=PP_ALIGN.RIGHT, fill=bg)
        _cell(table.cell(r, 8), D.br_dec(carteira), size=10.5, color=T.BODY,
              font=T.F_MONO, align=PP_ALIGN.RIGHT, fill=bg)
        table.rows[r].height = Inches(0.37)
    # total
    rr = nrows - 1
    mvals, caixa, carteira = total_row
    _cell(table.cell(rr, 0), "", fill=T.INK)
    _cell(table.cell(rr, 1), "TOTAL", size=11, bold=True, color=T.WHITE,
          font=T.F_HEAD, align=PP_ALIGN.LEFT, fill=T.INK)
    for j, mv in enumerate(mvals):
        _cell(table.cell(rr, 2 + j), D.br_int(mv), size=10.5, bold=True,
              color=T.WHITE, font=T.F_MONO, align=PP_ALIGN.RIGHT, fill=T.INK)
    _cell(table.cell(rr, 7), D.br_dec(caixa), size=10.5, bold=True, color=T.WHITE,
          font=T.F_MONO, align=PP_ALIGN.RIGHT, fill=T.INK)
    _cell(table.cell(rr, 8), D.br_dec(carteira), size=10.5, bold=True, color=T.WHITE,
          font=T.F_MONO, align=PP_ALIGN.RIGHT, fill=T.INK)
    table.rows[rr].height = Inches(0.44)
    _table_borders(table, ncols, nrows)
    # nota
    txt(s, LM, Inches(6.85), CW, Inches(0.3),
        [{"text": "Coluna MAI destacada · “–” indica produto ainda não lançado no mês · "
                  "Carteira em R$.", "size": 9, "color": T.MUTED, "font": T.F_BODY,
          "italic": True}])
    footer(s, page)


def _table_borders(table, ncols, nrows):
    """Bordas finas horizontais; remove verticais para leveza."""
    for r in range(nrows):
        for c in range(ncols):
            cell = table.cell(r, c)
            tcPr = cell._tc.get_or_add_tcPr()
            # remove existing borders
            for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
                for e in tcPr.findall(qn(tag)):
                    tcPr.remove(e)
            # linha inferior fina
            for tag, color, w in (('a:lnB', 'D4DEE1', 9525),):
                ln = tcPr.makeelement(qn(tag), {'w': str(w), 'cap': 'flat'})
                fill = ln.makeelement(qn('a:solidFill'), {})
                clr = fill.makeelement(qn('a:srgbClr'), {'val': color})
                fill.append(clr)
                ln.append(fill)
                tcPr.append(ln)


# ============================================================
# BUILD
# ============================================================
slide_capa()
slide_in_out()
slide_in_out_faturar()
slide_venda_producao()
slide_por_produto(1, "Venda IN por Produto", "Comercial · Maio 2026", T.GREEN, 5)
slide_por_produto(2, "Venda OUT por Produto", "Comercial · Maio 2026", T.AMBER, 6)
slide_por_produto(3, "Produção por Produto", "Produção · Maio 2026", T.BLUE, 7)
slide_estoque()
slide_mensal("Vendas IN — 2026 (Jan a Mai)", "Histórico · Venda IN",
             D.VENDAS_IN, D.VENDAS_IN_TOTAL, T.GREEN, 9, "IN")
slide_mensal("Vendas OUT — 2026 (Jan a Mai)", "Histórico · Venda OUT",
             D.VENDAS_OUT, D.VENDAS_OUT_TOTAL, T.AMBER, 10, "OUT")
slide_mensal("Produção — 2026 (Jan a Mai)", "Histórico · Produção",
             D.PRODUCAO_TAB, D.PRODUCAO_TOTAL, T.BLUE, 11, "PROD")

cp = prs.core_properties
cp.title = "Relatório de Vendas & Produção — Maio 2026"
cp.author = "CleanTabaco"
cp.subject = "Vendas IN/OUT, Produção e Estoque — Maio 2026"

import os
_OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                    "Apresentacao_Maio_2026.pptx")
prs.save(_OUT)
print("OK saved ->", _OUT)
