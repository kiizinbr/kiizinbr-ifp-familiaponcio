"""Sistema visual do deck — direção 'Premium / Dados Clínicos'.
Paleta restrita, tipografia clara, números em monoespaçada.
"""
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Dimensões do slide (16:9) ----
EMU_IN = 914400
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---- Cores ----
def C(hexstr):
    return RGBColor.from_string(hexstr)

WHITE   = C("FFFFFF")
CANVAS  = C("F5F8F9")   # fundo geral muito claro
PANEL   = C("FFFFFF")   # cartões
TINT    = C("EEF3F4")   # painel sutil
INK     = C("16242D")   # títulos
BODY    = C("3A4A52")   # texto corrente
MUTED   = C("7C8B92")   # texto secundário
HAIR    = C("E1E8EB")   # linhas finas
HAIR2   = C("D4DEE1")

BRAND   = C("0E7C86")   # teal — acento institucional
BRAND_D = C("0A5B62")   # teal escuro
BRAND_L = C("E2F0F1")   # teal claro (fundos)

# Cores semânticas de dados (consistentes em todo o deck)
GREEN   = C("2E9E62")   # Venda IN / Venda
GREEN_L = C("E4F3EA")
AMBER   = C("E08A2B")   # Venda OUT
AMBER_L = C("FBEEDD")
GOLD    = C("C29A2E")   # A faturar
GOLD_L  = C("F6EFD8")
BLUE    = C("2E6FA3")   # Produção
BLUE_L  = C("E3EDF5")
SLATE   = C("5B7079")   # estoque / neutro
SLATE_L = C("E9EEF0")

# Tons sequenciais para barras de produto (do mais forte ao mais suave)
def shades(base_hex, n):
    """Gera n tons do mesmo matiz, do cheio ao claro (mix com branco)."""
    base = tuple(int(base_hex[i:i+2], 16) for i in (0, 2, 4))
    out = []
    for k in range(n):
        # de 100% (k=0) até ~38% de intensidade
        f = 1.0 - (k / max(n - 1, 1)) * 0.62
        rgb = tuple(round(c * f + 255 * (1 - f)) for c in base)
        out.append(RGBColor(*rgb))
    return out

# ---- Tipografia ----
F_HEAD = "Segoe UI Semibold"
F_BODY = "Segoe UI"
F_LIGHT = "Segoe UI Light"
F_MONO = "Consolas"   # números / dados


# ---- Helpers de shape ----
def no_line(shape):
    shape.line.fill.background()


def no_shadow(shape):
    shape.shadow.inherit = False


def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    no_line(shape)
    no_shadow(shape)
    return shape


def rect(slide, l, t, w, h, color=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    no_shadow(sp)
    if color is not None:
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
    else:
        sp.fill.background()
    no_line(sp)
    return sp


def rounded(slide, l, t, w, h, color, radius=0.06):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    no_shadow(sp)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    no_line(sp)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def set_border(shape, color, width_pt=1.0):
    ln = shape.line
    ln.color.rgb = color
    ln.width = Pt(width_pt)


def soft_shadow(shape, blur=0.06, dist=0.04, alpha=78):
    """Sombra sutil e premium via XML."""
    spPr = shape.fill._xPr  # element with spPr semantics
    # find spPr
    el = shape._element.spPr
    # remove existing
    for tag in ("a:effectLst",):
        for e in el.findall(qn(tag)):
            el.remove(e)
    eff = el.makeelement(qn("a:effectLst"), {})
    sh = eff.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(blur * EMU_IN)),
        "dist": str(int(dist * EMU_IN)),
        "dir": "5400000",
        "rotWithShape": "0",
    })
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "1A2A33"})
    a = clr.makeelement(qn("a:alpha"), {"val": str(int((100 - alpha) * 1000 + 0))})
    # alpha: lower => more transparent. We want subtle -> alpha ~ 22000
    a.set("val", "20000")
    clr.append(a)
    sh.append(clr)
    eff.append(sh)
    el.append(eff)


def txt(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        wrap=True):
    """lines: lista de dicts {text, size, bold, color, font, spacing_after, ...}
       ou lista de listas de runs para múltiplos runs por parágrafo.
    """
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = line.get("align", align)
        if "space_before" in line:
            p.space_before = Pt(line["space_before"])
        if "space_after" in line:
            p.space_after = Pt(line["space_after"])
        if "line_spacing" in line:
            p.line_spacing = line["line_spacing"]
        runs = line.get("runs")
        if runs is None:
            runs = [line]
        for rdef in runs:
            r = p.add_run()
            r.text = rdef.get("text", "")
            f = r.font
            f.size = Pt(rdef.get("size", 14))
            f.bold = rdef.get("bold", False)
            f.name = rdef.get("font", F_BODY)
            f.color.rgb = rdef.get("color", BODY)
            if rdef.get("italic"):
                f.italic = True
            sp = rdef.get("spacing")
            if sp is not None:
                _set_char_spacing(r, sp)
    return tb


def _set_char_spacing(run, pts):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


# ---- Estrutura: cabeçalho / rodapé ----
def header(slide, kicker, title, accent=BRAND):
    """Cabeçalho consistente: kicker pequeno + título + régua de acento."""
    LM = Inches(0.62)
    # barra de acento vertical
    bar = rect(slide, LM, Inches(0.52), Inches(0.075), Inches(0.86), accent)
    bar.adjustments  # noop
    tx = LM + Inches(0.22)
    txt(slide, tx, Inches(0.46), Inches(11.4), Inches(0.34),
        [{"text": kicker.upper(), "size": 11.5, "bold": True, "color": accent,
          "font": F_HEAD, "spacing": 2.4}])
    txt(slide, tx, Inches(0.74), Inches(11.6), Inches(0.7),
        [{"text": title, "size": 25, "bold": True, "color": INK,
          "font": F_HEAD}])


def footer(slide, page_no, total=11):
    LM = Inches(0.62)
    y = Inches(7.06)
    rect(slide, LM, y, Inches(12.09), Pt(0.9), HAIR)
    txt(slide, LM, y + Inches(0.06), Inches(6), Inches(0.3),
        [{"runs": [
            {"text": "CLEANTABACO", "size": 9, "bold": True, "color": BRAND,
             "font": F_HEAD, "spacing": 1.5},
        ]}])
    txt(slide, Inches(5.5), y + Inches(0.06), Inches(3.3), Inches(0.3),
        [{"text": "RELATÓRIO MENSAL · MAIO 2026", "size": 9, "color": MUTED,
          "font": F_BODY, "align": PP_ALIGN.CENTER, "spacing": 1.0}],
        align=PP_ALIGN.CENTER)
    txt(slide, Inches(10.0), y + Inches(0.06), Inches(2.71), Inches(0.3),
        [{"text": f"{page_no:02d} / {total:02d}", "size": 9, "bold": True,
          "color": INK, "font": F_MONO, "align": PP_ALIGN.RIGHT}],
        align=PP_ALIGN.RIGHT)


def canvas_bg(slide, color=WHITE):
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)
